import PyPDF2
import io
import requests
from bs4 import BeautifulSoup

class DocumentExtractor:
    @staticmethod
    def extract_text_from_pdf(pdf_bytes: bytes) -> str:
        """Extracts text from a given PDF file in bytes."""
        try:
            pdf_file = io.BytesIO(pdf_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages[:5]:  # Limit to the first 5 pages
                text += page.extract_text()
            return text
        except Exception as e:
            print(f"An error occurred while extracting text: {e}")
            return ""

    @staticmethod
    def extract_text_from_pptx(pptx_bytes: bytes) -> str:
        """Extracts text from a given PPTX file in bytes."""
        try:
            from pptx import Presentation
            pptx_file = io.BytesIO(pptx_bytes)
            presentation = Presentation(pptx_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"An error occurred while extracting text: {e}")
            return ""
        
    @staticmethod
    def extract_text_from_webpage(url: str) -> str:
        """Extracts visible text from a given webpage URL."""
        try:

            response = requests.get(url, timeout=10)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, "html.parser")

            # Remove script and style elements
            for element in soup(["script", "style"]):
                element.decompose()

            text = soup.get_text(separator="\n")
        
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            return " ".join(lines)
        except Exception as e:
            print(f"An error occurred while extracting text from webpage: {e}")
            return ""
