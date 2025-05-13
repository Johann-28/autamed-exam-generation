import PyPDF2
import io

def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extracts text from a given PDF file in bytes."""
    try:
        pdf_file = io.BytesIO(pdf_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages[:5]:  # Limit to the first 5 pages
            text += page.extract_text()
        return text
    except Exception as e:
        print(f"An error occurred while extracting text: {e}")
        return ""

def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extracts text from a given PPTX file in bytes."""
    try:
        from pptx import Presentation
        pptx_file = io.BytesIO(pptx_bytes)
        presentation = Presentation(pptx_file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"An error occurred while extracting text: {e}")
        return ""
